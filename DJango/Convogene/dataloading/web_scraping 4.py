import pandas as pd
from langchain_community.document_loaders import AsyncHtmlLoader
from langchain_community.document_transformers import Html2TextTransformer
from langchain_text_splitters import RecursiveCharacterTextSplitter
from tqdm import tqdm
import tiktoken
import logging
from typing import Iterable
from langchain.schema import Document
import os
import io
import requests
import argparse
from pptx import Presentation
from pdfminer.high_level import extract_text
from fake_useragent import UserAgent
from dotenv import load_dotenv
import urllib.parse

# DATA PREPARATION

logging.basicConfig(level=logging.INFO, format='%(levelname)s: %(message)s')
logger = logging.getLogger(__name__)

load_dotenv()

class DataEngine:
    def __init__(self, input_file: str, output_dir: str, batch_size: int = 50):
        self.batch_size = batch_size
        self.input_file = input_file
        self.output_dir = output_dir
        self.ua = UserAgent()

    def get_links(self):
        if self.input_file.endswith(".csv"):
            df = pd.read_csv(self.input_file, encoding='ISO-8859-1')
            return df['doc_urls'].tolist()
        elif self.input_file.endswith(".txt"):
            with open(self.input_file, "r") as file:
                urls = file.readlines()
                urls = [url.strip() for url in urls if url.strip()]
                return urls

    def filter_urls(self, urls):
        filtered_urls = []
        remaining_urls = []
    
        try:
            for link in urls:
                parsed_url = urllib.parse.urlparse(link)
                path = parsed_url.path
                try:
                    headers = {'User-Agent': self.ua.random}
                    response = requests.get(link, headers=headers, verify=True)
                except requests.exceptions.SSLError:
                    logger.warning(f"Ignoring URL due to SSL verification error: {link}")
                    continue
                except requests.exceptions.RequestException as e:
                    logger.warning(f"Request error for URL {link}: {e}")
                    continue

                if response.status_code == 200 and not path.endswith((".mp4", ".mp3", ".pdf", ".pptx", ".ppt")):
                    filtered_urls.append(link)
                elif response.status_code == 200 and path.endswith((".pdf", ".pptx", ".ppt")):
                    remaining_urls.append(link)
                else:
                    logger.warning(f"Failed to fetch URL: {link}")

            logger.info(f"Number of filtered URLs: {len(filtered_urls)}")
            logger.info(f"Number of remaining URLs: {len(remaining_urls)}")
            
            return filtered_urls, remaining_urls
        
        except Exception as e:
            logger.error(f"Error processing some links: {e}")
            return [], []

    def pdf_pptx_text_extractor(self, remaining_urls):
        pdfdocs = []
        pptdocs = []
        
        for url in remaining_urls:
            if url.endswith('.pdf'):
                logger.info(f"Found PDF URL: {url}")
                try:
                    response = requests.get(url)
                    if response.status_code == 200:
                        text = extract_text(io.BytesIO(response.content))
                        pdfdocs.append(Document(page_content=text, metadata={'source': url}))
                    else:
                        logger.warning(f"Failed to fetch URL: {url}")
                except Exception as e:
                    logger.error(f"Error extracting text from PDF: {e}")

            elif url.endswith('.pptx'):
                logger.info(f"Found PPT URL: {url}")
                try:
                    response = requests.get(url)
                    if response.status_code == 200:
                        pptx_file = io.BytesIO(response.content)
                        presentation = Presentation(pptx_file)
                        text = ""
                        for slide in presentation.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text += shape.text + "\n"
                        pptdocs.append(Document(page_content=text, metadata={'source': url}))
                    else:
                        logger.warning(f"Failed to fetch URL: {url}")
                except Exception as e:
                    logger.error(f"Error extracting text from PowerPoint file: {e}")
        
        return pdfdocs, pptdocs

    def web_text_extractor(self, urls):
        filtered_urls, remaining_urls = self.filter_urls(urls)
        
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        
        async def get_docs(urls: Iterable[str]) -> Iterable[Document]:
            loader = AsyncHtmlLoader(concurrency=2, save_directory='html_files')
            docs = await loader.aget_docs(urls)
            return docs
        
        documents = loop.run_until_complete(get_docs(filtered_urls))
        
        transformer = Html2TextTransformer()
        documents = transformer.transform_documents(documents)
        
        pdfdocs, pptdocs = self.pdf_pptx_text_extractor(remaining_urls)
        documents.extend(pdfdocs)
        documents.extend(pptdocs)

        logger.info(f"Number of documents extracted: {len(documents)}")

        return documents

    def text_splitter(self, documents):
        encoding = tiktoken.encoding_for_model("gpt-3.5-turbo")

        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        documents = text_splitter.split_documents(documents)
        
        lengths = [len(encoding.encode(doc.page_content)) for doc in documents]
        logger.info(f"Number of documents after splitting: {len(documents)}")

        plt.hist(lengths, bins=50)
        plt.xlabel('Length of Text')
        plt.ylabel('Frequency')
        plt.title('Distribution of Text Lengths')
        plt.show()

        return documents

    def save_docs_to_csv(self, documents):
        output_file = os.path.join(self.output_dir, "processed_docs.csv")
        
        data = [{"page_content": doc.page_content, "metadata": doc.metadata} for doc in documents]
        df = pd.DataFrame(data)
        df.to_csv(output_file, index=False)
        
        logger.info(f"Documents saved to {output_file}")

    def process_documents(self):
        urls = self.get_links()
        documents = []
        
        for i in range(0, len(urls), self.batch_size):
            batch_urls = urls[i:i+self.batch_size]
            batch_docs = self.web_text_extractor(batch_urls)
            documents.extend(batch_docs)
        
        documents = self.text_splitter(documents)
        self.save_docs_to_csv(documents)

if __name__ == "__main__":
    # parser = argparse.ArgumentParser(description="Data Processor")
    # parser.add_argument("--input_file", type=str, required=True, help="Path to the input file containing URLs")
    # parser.add_argument("--output_dir", type=str, required=True, help="Directory to save the processed documents")
    # args = parser.parse_args()

    # processor = DataEngine(args.input_file, args.output_dir)
    processor.process_documents()
